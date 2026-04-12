from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("docs/reports/2026-04-13-positives-monday-strategy-deck-keynote-ready.pptx")
LOGO = Path("public/logos/positives-wordmark-dark.png")
HOME_SHOT = Path("tmp/slide-shots/home.png")
JOIN_SHOT = Path("tmp/slide-shots/join.png")
LOGIN_SHOT = Path("tmp/slide-shots/login.png")
WATCH_SHOT = Path("tmp/slide-shots/watch.png")

FONT_HEAD = "Avenir Next"
FONT_BODY = "Avenir Next"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

BG = RGBColor(248, 244, 238)
BG_SOFT = RGBColor(251, 248, 243)
WHITE = RGBColor(255, 255, 255)
CHARCOAL = RGBColor(18, 20, 23)
MUTED = RGBColor(104, 112, 122)
MUTED_2 = RGBColor(154, 160, 168)
BORDER = RGBColor(221, 215, 207)
STONE = RGBColor(239, 234, 226)
SAGE = RGBColor(111, 138, 122)
MOSS = RGBColor(78, 106, 89)
TEAL = RGBColor(46, 196, 182)
BLUE = RGBColor(47, 111, 237)
BLUE_DARK = RGBColor(36, 93, 208)
DARK = RGBColor(24, 28, 32)


def set_background(slide, rgb=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_shape(slide, shape_type, x, y, w, h, fill=None, line=None, radius=False, transparency=None):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else shape_type
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if transparency is not None:
            shape.fill.transparency = transparency
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=18,
    color=CHARCOAL,
    bold=False,
    font=FONT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
    margin=0.04,
    italic=False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    p.alignment = align
    return box


def add_lines(slide, x, y, w, lines, *, size=15, color=CHARCOAL, prefix="• ", gap=4):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3 + len(lines) * 0.36))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{prefix}{line}" if prefix else line
        p.space_after = Pt(gap)
        f = p.font
        f.name = FONT_BODY
        f.size = Pt(size)
        f.color.rgb = color
    return box


def add_logo(slide, x=0.84, y=0.48, w=1.25):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(x), Inches(y), width=Inches(w))


def add_orb(slide, x, y, w, h, color, transparency=0.82):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, y, w, h, fill=color, line=None, transparency=transparency)


def add_header(slide, title, subtitle=None):
    add_logo(slide, 0.82, 0.42, 1.18)
    add_orb(slide, 10.1, -0.2, 2.7, 2.0, TEAL, 0.9)
    add_orb(slide, 9.4, 0.15, 3.2, 2.35, BLUE, 0.92)
    add_text(slide, 0.82, 1.02, 7.6, 0.55, title, size=28, bold=True, font=FONT_HEAD)
    if subtitle:
        add_text(slide, 0.84, 1.48, 7.0, 0.28, subtitle, size=12.5, color=MUTED)
    rule = add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0.82, 1.8, 11.5, 0.02, fill=BORDER, line=None)
    rule.fill.solid()


def add_footer(slide, page):
    add_text(slide, 0.84, 6.98, 3.2, 0.18, "Positives • Monday strategy discussion", size=9.5, color=MUTED_2)
    add_text(slide, 11.8, 6.98, 0.28, 0.18, str(page), size=9.5, color=MUTED_2, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title=None, accent=None, fill=WHITE):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=BORDER, radius=True)
    if accent:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + 0.14, y + 0.13, w - 0.28, 0.07, fill=accent, line=None, radius=True)
    if title:
        add_text(slide, x + 0.18, y + 0.2, w - 0.36, 0.22, title, size=15, bold=True, font=FONT_HEAD)


def add_picture(slide, path, x, y, w):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def set_cell(cell, text, *, fill=WHITE, color=CHARCOAL, bold=False, size=14, font=FONT_BODY):
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    f = p.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color


def build_cover():
    slide = prs.slides.add_slide(blank)
    set_background(slide, BG_SOFT)
    add_orb(slide, 8.8, -0.2, 4.0, 3.2, BLUE, 0.8)
    add_orb(slide, 9.8, 0.7, 2.4, 2.0, TEAL, 0.76)
    add_orb(slide, -0.4, 5.9, 1.8, 1.5, SAGE, 0.9)
    add_logo(slide, 0.84, 0.56, 1.4)
    add_text(slide, 0.84, 1.58, 6.8, 0.98, "Positives Launch Planning", size=31, bold=True, font=FONT_HEAD)
    add_text(
        slide,
        0.86,
        2.72,
        6.9,
        0.46,
        "What we've built, what's changing, and what we need to decide",
        size=18,
        color=MUTED,
    )
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.88, 4.56, 5.8, 1.3, fill=WHITE, line=None, radius=True)
    add_text(
        slide,
        1.14,
        4.9,
        5.2,
        0.42,
        "A guided conversation deck for moving from the legacy stack to an owned Positives platform.",
        size=17,
    )
    add_text(slide, 0.88, 6.52, 4.3, 0.22, "Dr. Paul Jenkins + Ryan Bradshaw", size=11.5, color=MUTED)
    add_text(slide, 11.82, 6.92, 0.28, 0.18, "1", size=9.5, color=MUTED_2, align=PP_ALIGN.RIGHT)


def build_where_we_are_now():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "We've built a real platform")
    add_text(
        slide,
        0.92,
        2.15,
        5.25,
        1.05,
        "Positives is no longer just an idea. The core membership, billing, and admin foundation already exist.",
        size=24,
        bold=True,
        font=FONT_HEAD,
    )
    add_lines(
        slide,
        0.98,
        3.58,
        4.85,
        [
            "core membership experience exists",
            "billing foundation exists",
            "admin system exists",
            "launch decisions are now mostly business and migration decisions",
        ],
        size=15,
    )
    add_card(slide, 6.55, 2.04, 5.6, 4.92, "Live production proof", accent=SAGE)
    add_picture(slide, HOME_SHOT, 6.8, 2.42, 5.1)
    add_text(slide, 6.84, 6.48, 4.9, 0.18, "Live now at `positives.life`", size=11.5, color=MUTED)
    add_footer(slide, 2)


def build_what_exists():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "What's already built")
    items = [
        ("Stripe checkout and billing portal", TEAL),
        ("Daily practice member experience", SAGE),
        ("Library, notes, and account area", BLUE),
        ("Events and coaching surfaces", SAGE),
        ("Admin content management", TEAL),
        ("Affiliate / referral foundation", BLUE),
        ("Public funnel and marketing pages", SAGE),
    ]
    positions = [
        (0.92, 2.18), (4.28, 2.18), (7.64, 2.18),
        (0.92, 4.05), (4.28, 4.05), (7.64, 4.05),
        (2.6, 5.92),
    ]
    sizes = [(3.0, 1.42)] * 6 + [(8.0, 0.84)]
    for (label, accent), (x, y), (w, h) in zip(items, positions, sizes):
        add_card(slide, x, y, w, h, accent=accent)
        add_text(slide, x + 0.18, y + 0.46, w - 0.36, h - 0.44, label, size=14.2, bold=True if h > 1 else False, font=FONT_HEAD if h > 1 else FONT_BODY)
    add_footer(slide, 3)


def build_product_proof():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "The new member experience is real")
    add_card(slide, 0.84, 2.02, 5.85, 2.2, "Homepage", accent=TEAL)
    add_card(slide, 6.64, 2.02, 5.85, 2.2, "Join page", accent=BLUE)
    add_card(slide, 0.84, 4.48, 5.85, 2.18, "Watch page", accent=SAGE)
    add_card(slide, 6.64, 4.48, 5.85, 2.18, "Login", accent=BLUE)
    add_picture(slide, HOME_SHOT, 1.0, 2.42, 5.5)
    add_picture(slide, JOIN_SHOT, 6.8, 2.42, 5.5)
    add_picture(slide, WATCH_SHOT, 1.0, 4.88, 5.5)
    add_picture(slide, LOGIN_SHOT, 6.8, 4.88, 5.5)
    add_footer(slide, 4)


def build_why_this_matters():
    slide = prs.slides.add_slide(blank)
    set_background(slide, BG_SOFT)
    add_header(slide, "This changes what's possible")
    add_text(
        slide,
        0.92,
        2.24,
        10.6,
        0.88,
        "The future no longer has to be defined by the old systems.",
        size=30,
        bold=True,
        font=FONT_HEAD,
    )
    add_card(slide, 0.92, 4.15, 3.55, 1.48, accent=TEAL)
    add_card(slide, 4.9, 4.15, 3.55, 1.48, accent=BLUE)
    add_card(slide, 8.88, 4.15, 3.44, 1.48, accent=SAGE)
    add_text(slide, 1.16, 4.54, 3.0, 0.44, "New members can join through the new platform.", size=16)
    add_text(slide, 5.14, 4.54, 3.0, 0.44, "Billing can move to Stripe over time.", size=16)
    add_text(slide, 9.12, 4.54, 2.9, 0.44, "The team can operate content directly.", size=16)
    add_text(slide, 0.92, 6.25, 6.5, 0.24, "The launch conversation can stay focused on transition and decisions, not whether the product exists.", size=13, color=MUTED)
    add_footer(slide, 5)


def build_legacy_vs_future():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "From legacy stack to owned platform")
    add_card(slide, 0.92, 2.04, 4.78, 4.82, "Current world", accent=STONE, fill=WHITE)
    add_card(slide, 7.54, 2.04, 4.78, 4.82, "Future world", accent=SAGE, fill=WHITE)
    add_lines(slide, 1.18, 2.64, 3.9, [
        "Infusionsoft / Keap",
        "Authorize.net",
        "ClickFunnels",
        "Liquid Web",
        "WordPress + plugin stack",
    ], size=15)
    add_lines(slide, 7.8, 2.64, 3.95, [
        "Stripe",
        "ActiveCampaign Pro",
        "Postmark",
        "Vercel",
        "Supabase",
        "FirstPromoter",
        "Vimeo",
    ], size=15)
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(5.98), Inches(4.05), Inches(1.24), Inches(0.58))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLUE
    arrow.line.fill.background()
    add_text(slide, 6.02, 3.5, 1.2, 0.18, "owned", size=11, color=MOSS, bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)


def build_five_decisions():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "What we need to decide today")
    items = [
        "what each membership tier should promise",
        "whether Stripe is truly ready for live money",
        "how we move existing members to Stripe",
        "whether we approve the email-stack direction",
        "how much of the old offer / funnel stack we carry forward",
    ]
    for i, item in enumerate(items):
        y = 2.08 + i * 0.92
        add_card(slide, 1.28, y, 10.78, 0.68, accent=TEAL if i % 2 == 0 else BLUE)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 0.92, y + 0.08, 0.48, 0.48, fill=MOSS, line=None)
        add_text(slide, 0.92, y + 0.135, 0.48, 0.16, str(i + 1), size=12, color=WHITE, bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER)
        add_text(slide, 1.58, y + 0.18, 9.9, 0.22, item, size=15.5, font=FONT_BODY)
    add_footer(slide, 7)


def build_promise_tiers():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "Keep the promise simple and truthful")
    cols = [
        ("Membership", [
            "daily audio",
            "weekly reflection",
            "monthly theme",
            "full core library",
        ], TEAL),
        ("Membership + Events", [
            "everything in Membership",
            "live member events",
            "replay access",
        ], BLUE),
        ("Coaching Circle", [
            "everything in Membership + Events",
            "weekly live coaching",
            "coaching replays",
        ], SAGE),
    ]
    xs = [0.92, 4.4, 7.88]
    for (title, bullets, accent), x in zip(cols, xs):
        add_card(slide, x, 2.16, 3.06, 4.28, title, accent=accent)
        add_lines(slide, x + 0.2, 2.86, 2.55, bullets, size=14)
    add_card(slide, 0.92, 6.02, 10.02, 0.58, accent=SAGE)
    add_text(slide, 1.16, 6.18, 9.54, 0.18, "Do not promise extra certified-coach support unless it is truly staffed.", size=13.5, color=MOSS, bold=True, font=FONT_HEAD)
    add_footer(slide, 8)


def build_stripe_readiness():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "Before we launch, Stripe must be real")
    add_text(slide, 0.96, 2.16, 5.3, 0.78, "This is about live readiness, not just sandbox readiness.", size=24, bold=True, font=FONT_HEAD)
    checklist = [
        "bank account and payouts",
        "business details",
        "verification and tax requirements",
        "right long-term owner",
        "operational confidence before live money",
    ]
    for i, item in enumerate(checklist):
        y = 3.28 + i * 0.58
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 1.02, y + 0.05, 0.18, 0.18, fill=SAGE, line=None)
        add_text(slide, 1.34, y, 4.65, 0.18, item, size=15)
    add_card(slide, 7.0, 2.2, 5.18, 3.8, "Decision lens", accent=BLUE)
    add_lines(slide, 7.28, 2.9, 4.5, [
        "If anything material is still unresolved, we should treat that as a launch gating task.",
        "The question is whether Positives is ready to accept live money with confidence and clarity.",
    ], size=14, prefix="")
    add_footer(slide, 9)


def build_member_migration():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "Move existing members calmly, not all at once")
    steps = [
        "announcement phase",
        "new members to Stripe now",
        "try payment-data migration first",
        "guided payment-update flows if needed",
        "do not shut down the old system too early",
    ]
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 1.28, 4.1, 9.9, 0.03, fill=BORDER, line=None)
    for i, step in enumerate(steps):
        x = 0.92 + i * 2.25
        y = 2.7 if i % 2 == 0 else 3.28
        add_card(slide, x, y, 1.82, 1.24, accent=TEAL if i in (0, 1) else BLUE)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x + 0.12, y + 0.12, 0.28, 0.28, fill=MOSS, line=None)
        add_text(slide, x + 0.12, y + 0.15, 0.28, 0.14, str(i + 1), size=10, color=WHITE, bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.14, y + 0.48, 1.48, 0.48, step, size=12.2)
    add_footer(slide, 10)


def build_pricing_transition():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "Simplify pricing over time")
    add_text(slide, 0.96, 2.22, 6.3, 0.44, "Current legacy cohorts include prices like 67, 47, and 07.", size=22, bold=True, font=FONT_HEAD)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.18, 4.18, 2.0, 1.14, fill=WHITE, line=BORDER, radius=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 5.22, 4.18, 2.0, 1.14, fill=WHITE, line=BORDER, radius=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 9.26, 4.18, 2.0, 1.14, fill=WHITE, line=BORDER, radius=True)
    add_text(slide, 1.46, 4.55, 1.4, 0.22, "67 • 47 • 07", size=21, bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER)
    add_text(slide, 5.58, 4.55, 1.3, 0.22, "transition", size=19, color=MOSS, bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER)
    add_text(slide, 9.9, 4.55, 0.72, 0.22, "37", size=26, color=BLUE, bold=True, font=FONT_HEAD, align=PP_ALIGN.CENTER)
    for x in (3.58, 7.62):
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(4.52), Inches(0.84), Inches(0.32))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = SAGE
        arrow.line.fill.background()
    add_card(slide, 0.96, 5.96, 10.9, 0.58, accent=SAGE)
    add_text(slide, 1.18, 6.12, 10.4, 0.18, "Recommendation: move legacy cohorts toward the current 37 plan, with exceptions only when intentionally approved.", size=13.5, color=MOSS, bold=True, font=FONT_HEAD)
    add_footer(slide, 11)


def build_rebilling_waves():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "Phased rebilling protects revenue")
    waves = [
        ("1", "new members already on Stripe", TEAL),
        ("2", "simple monthly members", TEAL),
        ("3", "legacy monthly cohorts", BLUE),
        ("4", "annual members", BLUE),
        ("5", "special cases and edge cases", SAGE),
    ]
    for i, (num, label, accent) in enumerate(waves):
        x = 1.0 + i * 2.2
        y = 5.7 - i * 0.42
        h = 0.78 + i * 0.42
        add_card(slide, x, y, 1.8, h, accent=accent)
        add_text(slide, x + 0.16, y + 0.18, 0.3, 0.16, num, size=12, color=accent, bold=True, font=FONT_HEAD)
        add_text(slide, x + 0.16, y + 0.42, 1.42, h - 0.52, label, size=12)
    add_card(slide, 0.92, 2.18, 11.38, 0.56, accent=SAGE)
    add_text(slide, 1.18, 2.34, 10.9, 0.18, "Do not move the next wave until the last one is stable.", size=13.8, color=MOSS, bold=True, font=FONT_HEAD)
    add_footer(slide, 12)


def build_email_recommendation():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "A cleaner email system")
    blocks = [
        (1.18, 2.72, 2.45, "App", "triggers events"),
        (4.42, 2.72, 3.15, "ActiveCampaign Pro", "automation brain"),
        (8.38, 2.72, 2.45, "Postmark", "transactional delivery"),
    ]
    colors = [TEAL, BLUE, SAGE]
    for (x, y, w, title, desc), accent in zip(blocks, colors):
        add_card(slide, x, y, w, 1.4, accent=accent)
        add_text(slide, x + 0.18, y + 0.36, w - 0.36, 0.22, title, size=16, bold=True, font=FONT_HEAD)
        add_text(slide, x + 0.18, y + 0.72, w - 0.36, 0.2, desc, size=13, color=MUTED)
    for x in (3.72, 7.76):
        arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(3.15), Inches(0.42), Inches(0.28))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = BLUE
        arrow.line.fill.background()
    add_lines(slide, 1.2, 5.05, 10.4, [
        "ActiveCampaign becomes the automation brain",
        "Postmark becomes the transactional delivery layer",
        "essential account emails still send",
        "marketing emails stay unsubscribable",
    ], size=14)
    add_footer(slide, 13)


def build_rebuild_first():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "We do not need to move everything")
    panels = [
        ("Now", [
            "membership-critical experience",
            "highest-value supporting offers",
            "priority funnels",
        ], TEAL, 0.92),
        ("Later", [
            "long-tail courses",
            "old campaigns",
            "lower-value legacy funnel sprawl",
        ], BLUE, 4.46),
        ("Maybe never", [
            "anything that no longer fits the future Positives direction",
        ], SAGE, 8.0),
    ]
    widths = [3.0, 3.0, 4.32]
    for (title, bullets, accent, x), w in zip(panels, widths):
        add_card(slide, x, 2.2, w, 4.18, title, accent=accent)
        add_lines(slide, x + 0.18, 2.94, w - 0.36, bullets, size=14)
    add_footer(slide, 14)


def build_cost_story():
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header(slide, "The overlap is temporary. The simplification is the point.")
    add_card(slide, 0.92, 2.12, 3.48, 3.8, "Target Positives stack", accent=TEAL)
    add_card(slide, 4.9, 2.12, 3.48, 3.8, "Legacy stack to retire", accent=STONE)
    add_card(slide, 8.88, 2.12, 3.44, 3.8, "Overlap period", accent=BLUE)
    add_lines(slide, 1.18, 2.82, 2.9, [
        "ActiveCampaign Pro",
        "Postmark",
        "Vercel",
        "Supabase",
        "Stripe",
        "Vimeo",
    ], size=13.5)
    add_lines(slide, 5.16, 2.82, 2.9, [
        "Keap",
        "ClickFunnels",
        "Liquid Web",
        "legacy plugin stack",
    ], size=13.5)
    add_lines(slide, 9.16, 2.82, 2.72, [
        "costs more for a while",
        "temporary dual-system period",
        "cleaner operating model afterward",
    ], size=13.5)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.92, 6.14, 11.4, 0.58, fill=WHITE, line=None, radius=True)
    add_text(slide, 1.16, 6.3, 10.9, 0.18, "Long-term goal: a cleaner, lower-friction system with less legacy dependence.", size=13.5, color=MOSS, bold=True, font=FONT_HEAD)
    add_footer(slide, 15)


def build_outcome():
    slide = prs.slides.add_slide(blank)
    set_background(slide, BG_SOFT)
    add_orb(slide, 9.8, -0.15, 2.6, 2.0, BLUE, 0.88)
    add_orb(slide, 10.4, 0.56, 1.9, 1.6, TEAL, 0.8)
    add_logo(slide, 0.82, 0.45, 1.18)
    add_text(slide, 0.82, 1.02, 7.0, 0.58, "The outcome we want today", size=29, bold=True, font=FONT_HEAD)
    outcomes = [
        "approved membership promise",
        "clear Stripe readiness path",
        "phased rebilling approved",
        "ActiveCampaign Pro + Postmark approved",
        "membership-first migration scope approved",
        "owners for follow-up",
    ]
    for i, item in enumerate(outcomes):
        y = 2.18 + i * 0.66
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 1.0, y + 0.03, 0.2, 0.2, fill=SAGE, line=None)
        add_text(slide, 1.36, y, 5.6, 0.18, item, size=16)
    add_card(slide, 7.12, 2.16, 5.0, 3.86, "Meeting goal", accent=TEAL)
    add_text(
        slide,
        7.38,
        2.9,
        4.42,
        1.0,
        "Leave with real decisions, a clear transition path, and named next owners.",
        size=24,
        bold=True,
        font=FONT_HEAD,
    )
    add_text(slide, 7.38, 5.15, 4.2, 0.28, "Support the conversation. Do not try to narrate every detail.", size=12.5, color=MUTED)
    add_text(slide, 11.82, 6.92, 0.28, 0.18, "16", size=9.5, color=MUTED_2, align=PP_ALIGN.RIGHT)


def main():
    build_cover()
    build_where_we_are_now()
    build_what_exists()
    build_product_proof()
    build_why_this_matters()
    build_legacy_vs_future()
    build_five_decisions()
    build_promise_tiers()
    build_stripe_readiness()
    build_member_migration()
    build_pricing_transition()
    build_rebilling_waves()
    build_email_recommendation()
    build_rebuild_first()
    build_cost_story()
    build_outcome()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    main()
